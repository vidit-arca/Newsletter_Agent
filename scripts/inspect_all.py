from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx'

prs = Presentation(pptx_path)

def extract_text_from_shape(shape, depth=0):
    indent = "  " * depth
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}Group Shape:")
        for s in shape.shapes:
            extract_text_from_shape(s, depth+1)
    elif shape.has_text_frame:
        print(f"{indent}TextFrame: '{shape.text.replace(chr(10), ' ').replace(chr(11), ' ')}'")
    elif shape.has_table:
        print(f"{indent}Table:")
        for row in shape.table.rows:
            row_data = [cell.text.replace(chr(10), ' ').replace(chr(11), ' ') for cell in row.cells]
            print(f"{indent}  Row: {row_data}")
    else:
        pass # Ignore pictures for now

for i, slide in enumerate(prs.slides):
    print(f"\n--- SLIDE {i+1} ---")
    for j, shape in enumerate(slide.shapes):
        extract_text_from_shape(shape)
