from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

pptx_path = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx'

prs = Presentation(pptx_path)

def extract_text_from_shape(shape, depth=0):
    indent = "  " * depth
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        print(f"{indent}Group Shape:")
        for s in shape.shapes:
            extract_text_from_shape(s, depth+1)
    elif shape.has_text_frame:
        print(f"{indent}TextFrame: '{shape.text.replace(chr(10), ' ').replace(chr(11), ' ')}'")
    elif shape.has_table:
        print(f"{indent}Table:")
        for row in shape.table.rows:
            row_data = [cell.text.replace(chr(10), ' ').replace(chr(11), ' ') for cell in row.cells]
            print(f"{indent}  Row: {row_data}")
    else:
        print(f"{indent}Other Shape: {shape.shape_type}")

print("--- DETAILED SLIDE 5 (Index 4) ---")
slide = prs.slides[4]
for j, shape in enumerate(slide.shapes):
    print(f"\nShape {j}:")
    extract_text_from_shape(shape)
