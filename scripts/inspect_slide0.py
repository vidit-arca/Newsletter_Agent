from pptx import Presentation
prs = Presentation('/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx')
slide = prs.slides[0]
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        print(f"Shape {i}: {shape.text!r}")
