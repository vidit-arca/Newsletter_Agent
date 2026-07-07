from pptx import Presentation
from pptx.util import Inches

prs = Presentation('/Users/apple/Desktop/Akshayam/Newsletter_Agent/output.pptx')
print(f"Slide Height: {prs.slide_height.inches} inches")
print(f"Slide Width: {prs.slide_width.inches} inches")

print("\n--- Analysing text bounds on Slide 5 (First content slide) ---")
slide = prs.slides[4]
for j, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        text = shape.text.replace("\n", " ")[:30]
        top = shape.top.inches if hasattr(shape, "top") else "N/A"
        height = shape.height.inches if hasattr(shape, "height") else "N/A"
        bottom = top + height if top != "N/A" and height != "N/A" else "N/A"
        print(f"Shape {j} '{text}...': top={top}, height={height}, bottom={bottom}")
