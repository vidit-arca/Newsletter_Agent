import copy
import io
from pptx import Presentation

prs = Presentation('/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx')
template_slide = prs.slides[9]
new_slide = prs.slides.add_slide(template_slide.slide_layout)

bg = template_slide.element.xpath('./p:cSld/p:bg')
if bg:
    new_slide.element.cSld.insert(0, copy.deepcopy(bg[0]))

for shp in template_slide.shapes:
    if shp.shape_type == 13: # PICTURE
        image_stream = io.BytesIO(shp.image.blob)
        new_slide.shapes.add_picture(image_stream, shp.left, shp.top, shp.width, shp.height)
    else:
        el = shp.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')

prs.save('/Users/apple/Desktop/Akshayam/Newsletter_Agent/test_duplicate_bg.pptx')
print('Done copying slide with background and image!')
