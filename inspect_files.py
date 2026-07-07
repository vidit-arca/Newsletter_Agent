import pandas as pd
from pptx import Presentation

excel_path = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/data/SEBI (2).xlsx'
pptx_path = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx'

print("--- EXCEL DATA ---")
try:
    df = pd.read_excel(excel_path)
    print("Columns:", df.columns.tolist())
    print("\nFirst few rows:")
    print(df.head(2).to_string())
except Exception as e:
    print("Error reading Excel:", e)

print("\n--- PPTX DATA ---")
try:
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"\nSlide {i+1}:")
        for j, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                text = shape.text[:50].replace('\n', ' ')
                print(f"  Shape {j}: Text='{text}'")
            else:
                print(f"  Shape {j}: Type={shape.shape_type}")
except Exception as e:
    print("Error reading PPTX:", e)
