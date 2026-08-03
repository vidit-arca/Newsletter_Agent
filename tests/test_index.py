import os
import re
import glob
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

data_dir = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/data/'
excel_files = glob.glob(os.path.join(data_dir, '*.xlsx'))
sources = []
for f in excel_files:
    base = os.path.basename(f)
    name = re.sub(r'\s*\(\d+\)\.xlsx$', '', base)
    name = re.sub(r'\.xlsx$', '', name)
    sources.append(name)
if len(sources) > 1:
    sources_str = ", ".join(sources[:-1]) + " and " + sources[-1]
elif sources:
    sources_str = sources[0]
else:
    sources_str = "SEBI"

dfs = {}
for file in excel_files:
    file_dfs = pd.read_excel(file, sheet_name=None)
    for sheet_name, df_sheet in file_dfs.items():
        if sheet_name in dfs:
            dfs[sheet_name] = pd.concat([dfs[sheet_name], df_sheet], ignore_index=True)
        else:
            dfs[sheet_name] = df_sheet

prs = Presentation('Template/Pravartiya - Template (1).pptx')
index_slide = prs.slides[1]
txBox = index_slide.shapes.add_textbox(Inches(0.5), Inches(4.5), Inches(7.5), Inches(5.0))
tf = txBox.text_frame
tf.word_wrap = True

first_para = True
for sheet_name, df_sheet in dfs.items():
    if df_sheet.empty:
        continue
    if first_para:
        p = tf.paragraphs[0]
        first_para = False
    else:
        p = tf.add_paragraph()
    p.text = f"\n{sheet_name.upper()}:"
    p.font.bold = True
    p.font.underline = True
    p.font.size = Pt(12)
    for _, row in df_sheet.iterrows():
        title = str(row['Title']).strip() if 'Title' in row and pd.notna(row['Title']) else ""
        if title:
            bullet_p = tf.add_paragraph()
            bullet_p.text = f"➢ {title}"
            bullet_p.font.size = Pt(12)

prs.save('test_output.pptx')
print("Done!")
