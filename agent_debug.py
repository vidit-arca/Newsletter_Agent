import re
import os
import glob
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─────────────────────────────────────────────────────────────────────────────
# Helper: parse "Regulation Number / Gist of amendment" blocks from a summary
# ─────────────────────────────────────────────────────────────────────────────
def parse_regulation_blocks(summary_text: str) -> dict:
    """
    Parses a regulation summary string and returns an ordered dict:
        { regulation_number: [gist1, gist2, ...], ... }

    Each block in the text looks like:
        Regulation Number: 15
        Footer Number: 91
        Gist of amendment: Omitted: ...
        Existing provisions: ...
    """
    # Split on double-newline before a "Regulation Number:" line
    blocks = re.split(r'\n\n(?=Regulation Number:)', summary_text)

    # Preserve insertion order (Python 3.7+)
    reg_dict: dict = {}

    for block in blocks:
        reg_match = re.search(r'Regulation Number:\s*(.+)', block)
        gist_match = re.search(r'Gist of amendment:\s*(.+?)(?=\nExisting provisions|\nRegulation Number:|$)',
                                block, re.DOTALL)
        if reg_match:
            reg_num = reg_match.group(1).strip()
            gist = gist_match.group(1).strip().replace('\n', ' ') if gist_match else ''
            reg_dict.setdefault(reg_num, [])
            if gist:
                reg_dict[reg_num].append(gist)

    return reg_dict


# ─────────────────────────────────────────────────────────────────────────────
# Helper: add a styled table to a slide for regulations with >3 amendments
# ─────────────────────────────────────────────────────────────────────────────
# Characters that fit in the Amendment column at font size 9 (~5.8 in wide)
_AMEND_COL_CHARS = 95
_HEADER_H        = 0.30   # inches for the header row
_MIN_ROW_H       = 0.28   # inches for a single-line data row
# Minimum space needed on a slide before we bother starting another chunk.
# Prevents orphan tables of 1-2 rows at the bottom of a slide.
_MIN_CHUNK_SPACE = 2.5


def _estimate_row_height(gist: str, chars_per_line: int = _AMEND_COL_CHARS) -> float:
    """Estimate the height of one table data row based on gist text length."""
    lines = max(1, -(-len(gist) // chars_per_line))   # ceiling division
    return max(_MIN_ROW_H, lines * 0.18 + 0.10)


def add_regulation_table(slide, table_rows: list, left_in: float, top_in: float,
                          width_in: float = 7.2, row_offset: int = 0) -> float:
    """
    Adds a PPTX table with columns [Regulation Number | Amendment].
    table_rows  – list of (reg_num, gist) tuples (header excluded).
    row_offset  – used to keep alternating stripe colour consistent across chunks.
    Returns the actual height in inches consumed by the table.
    """
    num_rows = len(table_rows) + 1   # +1 for header
    num_cols = 2

    # Pre-compute per-row heights so the table is sized accurately
    row_heights_in = [_HEADER_H] + [_estimate_row_height(g) for _, g in table_rows]
    total_height_in = sum(row_heights_in)

    from pptx.util import Emu as _Emu
    row_heights_emu = [int(h * 914400) for h in row_heights_in]

    table = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(left_in), Inches(top_in),
        Inches(width_in), int(total_height_in * 914400)
    ).table

    # Apply per-row heights
    for i, h_emu in enumerate(row_heights_emu):
        table.rows[i].height = h_emu

    # Column widths
    table.columns[0].width = Inches(1.4)
    table.columns[1].width = Inches(width_in - 1.4)

    HEADER_BG = RGBColor(0, 51, 153)    # dark blue
    ALT_BG    = RGBColor(235, 241, 255) # light blue stripe
    WHITE     = RGBColor(255, 255, 255)
    DARK_TEXT = RGBColor(30, 30, 30)

    headers = ["Regulation Number", "Amendment"]
    for col_idx, header_text in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.text = header_text
        cell.fill.solid()
        cell.fill.fore_color.rgb = HEADER_BG
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.size = Pt(10)
        para.font.color.rgb = WHITE
        para.alignment = PP_ALIGN.CENTER

    for row_idx, (reg_num, gist) in enumerate(table_rows, start=1):
        is_alt = ((row_idx + row_offset) % 2 == 0)
        for col_idx, cell_text in enumerate([reg_num, gist]):
            cell = table.cell(row_idx, col_idx)
            cell.text = cell_text
            cell.fill.solid()
            cell.fill.fore_color.rgb = ALT_BG if is_alt else WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(9)
            para.font.color.rgb = DARK_TEXT
            para.word_wrap = True
            if col_idx == 0:
                para.alignment = PP_ALIGN.CENTER
                para.font.bold = True

    return total_height_in


# ─────────────────────────────────────────────────────────────────────────────
# Main newsletter builder
# ─────────────────────────────────────────────────────────────────────────────
def create_newsletter():
    data_dir      = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/data/'
    template_path = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/Template/Pravartiya - Template (1).pptx'
    output_path   = '/Users/apple/Desktop/Akshayam/Newsletter_Agent/output.pptx'

    excel_files = glob.glob(os.path.join(data_dir, '*.xlsx'))
    
    sources = []
    for f in excel_files:
        base = os.path.basename(f)
        name = re.sub(r'\s*\(\d+\)\.xlsx$', '', base)
        name = re.sub(r'\.xlsx$', '', name)
        sources.append(name)
    if len(sources) > 1:
        sources_str = ", ".join(sources[:-1]) + " and " + sources[-1]
    elif sources:
        sources_str = sources[0]
    else:
        sources_str = "SEBI"

    dfs = {}
    file_grouped_dfs = {}
    for file in excel_files:
        base = os.path.basename(file)
        name = re.sub(r'\s*\(\d+\)\.xlsx$', '', base)
        name = re.sub(r'\.xlsx$', '', name)
        
        if name not in file_grouped_dfs:
            file_grouped_dfs[name] = {}
            
        file_dfs = pd.read_excel(file, sheet_name=None)
        for sheet_name, df_sheet in file_dfs.items():
            def is_valid_row(row):
                for col in ['Summary', 'combined_summary', 'updated_combined_summary']:
                    if col in row and pd.notna(row[col]):
                        s = str(row[col]).lower().replace(' ', '')
                        if 'pdfignore' in s or 'pdfingore' in s or 'pdftoignore' in s or 'pdftoingore' in s:
                            return False
                return True

            if not df_sheet.empty:
                mask = df_sheet.apply(is_valid_row, axis=1)
                df_sheet = df_sheet[mask]
            
            if sheet_name in dfs:
                dfs[sheet_name] = pd.concat([dfs[sheet_name], df_sheet], ignore_index=True)
            else:
                dfs[sheet_name] = df_sheet
                
            if sheet_name in file_grouped_dfs[name]:
                file_grouped_dfs[name][sheet_name] = pd.concat([file_grouped_dfs[name][sheet_name], df_sheet], ignore_index=True)
            else:
                file_grouped_dfs[name][sheet_name] = df_sheet

    if dfs:
        all_df = pd.concat(dfs.values(), ignore_index=True)
    else:
        all_df = pd.DataFrame()

    if not all_df.empty:
        target_month = str(all_df.iloc[0]['Month'])
        target_year  = str(all_df.iloc[0]['Year'])
    else:
        target_month = "Month"
        target_year  = "Year"

    prs = Presentation(template_path)

    for i in range(min(3, len(prs.slides))):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text
                if "____" in text:
                    text = re.sub(r'_{4,}', target_month, text)
                if "from SEBI issued" in text:
                    text = text.replace("from SEBI issued", f"from {sources_str} issued")
                shape.text = text

    index_slide_idx = 1
    current_index_slide = prs.slides[index_slide_idx]
    
    current_y = 4.5
    txBox = current_index_slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(7.2), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    first_para = True
    for file_name, sheets_dict in file_grouped_dfs.items():
        has_data = any(not df_sheet.empty for df_sheet in sheets_dict.values())
        if not has_data:
            continue
            
        if current_y + 0.4 > 9.8:
            if index_slide_idx == 1:
                index_slide_idx = 2
                current_index_slide = prs.slides[index_slide_idx]
                current_y = 2.0
                txBox = current_index_slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(7.2), Inches(8.0))
                tf = txBox.text_frame
                tf.word_wrap = True
                first_para = True
                
        # Domain header intentionally removed as requested
        
        for sheet_name, df_sheet in sheets_dict.items():
            if df_sheet.empty:
                continue
                
            # Calculate dynamic space needed based on character length of titles
            estimated_needed = 0.35 # Header
            for _, row in df_sheet.iterrows():
                title = str(row['Title']).strip().replace('**', '') if 'Title' in row and pd.notna(row['Title']) else ""
                if title:
                    lines = max(1, len(title) // 80 + 1)
                    estimated_needed += (lines * 0.2) + 0.1
            
            if current_y + estimated_needed > 9.8:
                if index_slide_idx == 1:
                    index_slide_idx = 2
                    current_index_slide = prs.slides[index_slide_idx]
                    current_y = 2.0
                    txBox = current_index_slide.shapes.add_textbox(Inches(0.5), Inches(current_y), Inches(7.2), Inches(8.0))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    first_para = True
                    
            if first_para:
                p = tf.paragraphs[0]
                first_para = False
            else:
                p = tf.add_paragraph()
                
            p.text = f"\n{sheet_name.upper()}:"
            p.font.bold = True
            p.font.underline = True
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(0, 32, 96)
            current_y += 0.35
            
            for _, row in df_sheet.iterrows():
                title = str(row['Title']).strip().replace('**', '') if 'Title' in row and pd.notna(row['Title']) else ""
                if title:
                    bullet_p = tf.add_paragraph()
                    bullet_p.text = f"➢ {title}"
                    bullet_p.font.size = Pt(12)
                    bullet_p.space_after = Pt(6)
                    # Accurate increment per title accounting for wrapping
                    lines = max(1, len(title) // 80 + 1)
                    current_y += (lines * 0.2) + 0.1



    content_slides = [prs.slides[i] for i in range(4, 10)]
    slide_index    = 0
    current_slide  = content_slides[slide_index]

    start_top  = 2.0
    max_height = 10.5
    current_y  = start_top

    def ensure_slide_space(needed: float) -> bool:
        """Advance to next slide if needed. Duplicates the last slide if we run out."""
        nonlocal slide_index, current_slide, current_y
        if current_y > max_height - needed:
            slide_index += 1
            if slide_index < len(content_slides):
                current_slide = content_slides[slide_index]
                current_y = start_top
            else:
                import copy
                template_slide = content_slides[-1]
                new_slide = prs.slides.add_slide(prs.slide_layouts[6])
                for shp in template_slide.shapes:
                    el = shp.element
                    newel = copy.deepcopy(el)
                    new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
                
                # Move to before the end slide
                sldIdLst = prs.slides._sldIdLst
                new_sldId = sldIdLst[-1]
                sldIdLst.remove(new_sldId)
                sldIdLst.insert(-1, new_sldId)
                
                content_slides.append(new_slide)
                current_slide = new_slide
                current_y = start_top
        return True

    for sheet_name, df_sheet in dfs.items():
        if df_sheet.empty:
            continue

        if not ensure_slide_space(0.8):
            break

        # ── Section header ────────────────────────────────────────────────
        txBox = current_slide.shapes.add_textbox(
            Inches(0.5), Inches(current_y), Inches(7.2), Inches(0.5))
        tf = txBox.text_frame
        p  = tf.paragraphs[0]
        p.text = f"➤ {sheet_name.upper()}"
        p.font.bold  = True
        p.font.size  = Pt(16)
        p.font.color.rgb = RGBColor(0, 51, 153)
        current_y += 0.5

        # ── Determine the correct summary column ─────────────────────────
        is_regulations = (sheet_name.strip().lower() == 'regulations')

        for index, row in df_sheet.iterrows():
            title   = str(row.get('Title', '')).replace('**', '')
            pdf_url = str(row.get('PDF_URL', ''))

            # Pick right summary column
            if is_regulations:
                summary = str(row.get('updated_combined_summary',
                              row.get('combined_summary', ''))).replace('**', '')
            else:
                summary = str(row.get('Summary', '')).replace('**', '')

            # ── Regulations sheet: group by Regulation Number ─────────────
            if is_regulations and summary.strip():
                # Render title first
                if not ensure_slide_space(0.5):
                    break

                txBox = current_slide.shapes.add_textbox(
                    Inches(0.5), Inches(current_y), Inches(7.2), Inches(0.5))
                tf  = txBox.text_frame
                tf.word_wrap = True
                p   = tf.paragraphs[0]
                run = p.add_run()
                run.text = title
                run.font.bold  = True
                run.font.size  = Pt(12)
                if pdf_url and pdf_url != 'nan':
                    run.hyperlink.address    = pdf_url
                    run.font.color.rgb       = RGBColor(0, 102, 204)
                    run.font.underline       = True
                else:
                    run.font.color.rgb = RGBColor(0, 0, 0)
                p.space_after = Pt(6)
                txBox.height  = Inches(0.30)
                current_y    += 0.30 + 0.15

                # Parse regulation blocks
                reg_dict = parse_regulation_blocks(summary)

                # All regulation entries go into one unified table
                table_rows: list = []   # (reg_num, gist) pairs for the table

                for reg_num, gists in reg_dict.items():
                    for gist in gists:
                        table_rows.append((reg_num, gist))

                # ── Render table in slide-fitting chunks ──────────────────
                if table_rows:
                    remaining = list(table_rows)
                    row_offset = 0  # keeps stripe colour consistent across chunks

                    while remaining:
                        # Available space on current slide (keep 0.5 in margin)
                        available = max_height - current_y - 0.5

                        # Advance slide if:
                        #   a) Not enough space for even header + 1 row, OR
                        #   b) Space is below _MIN_CHUNK_SPACE (avoids orphan mini-tables)
                        if available < _MIN_CHUNK_SPACE:
                            ensure_slide_space(_MIN_CHUNK_SPACE)
                            available = max_height - current_y - 0.5

                        # Fit as many rows as possible into available space
                        chunk = []
                        used = _HEADER_H
                        for row in remaining:
                            rh = _estimate_row_height(row[1])
                            if used + rh > available and chunk:
                                break   # this row spills over — save for next slide
                            chunk.append(row)
                            used += rh

                        if not chunk:
                            # Single row is taller than available — force it in anyway
                            chunk = [remaining[0]]

                        actual_h = add_regulation_table(
                            current_slide, chunk,
                            left_in=0.5, top_in=current_y,
                            width_in=7.2, row_offset=row_offset)
                        current_y += actual_h + 0.20
                        row_offset += len(chunk)
                        remaining = remaining[len(chunk):]

                current_y += 0.15  # gap after entry

            # ── All other sheets: original sequential rendering ────────────
            else:
                paragraphs_to_print = []
                paragraphs_to_print.append(('title', title, pdf_url)); print(f'Processing title: {title[:30]}')
                for p_text in summary.split('\n'):
                    if p_text.strip():
                        paragraphs_to_print.append(('summary', p_text.strip(), None))

                while paragraphs_to_print:
                    if not ensure_slide_space(0.5):
                        break
                    if slide_index >= len(content_slides):
                        print('BREAKING because slide_index >= len(content_slides)')
                        break

                    txBox = current_slide.shapes.add_textbox(
                        Inches(0.5), Inches(current_y), Inches(7.2), Inches(0.5))
                    tf  = txBox.text_frame
                    tf.word_wrap = True

                    box_height  = 0.0
                    first_para  = True

                    while paragraphs_to_print:
                        ptype, text, url = paragraphs_to_print[0]

                        text_clean = text.strip().replace('"', '')
                        is_sub  = False
                        is_list = False

                        if ptype != 'title':
                            if 0 < len(text_clean) < 80:
                                if text_clean.endswith(':') or text_clean.lower() in [
                                    "background and facts", "query",
                                    "response from sebi", "conclusion",
                                    "regulation reference given by sebi"]:
                                    is_sub = True
                            if text.strip().startswith('-') or text.strip().startswith('•'):
                                is_list = True

                        # Dynamic padding
                        if ptype == 'title':
                            para_padding = 0.15
                        elif is_sub:
                            para_padding = 0.25
                        elif is_list:
                            para_padding = 0.05
                        else:
                            para_padding = 0.12

                        chars_per_line = 90
                        lines      = int(len(text) / chars_per_line) + 1
                        para_height = lines * 0.18 + para_padding

                        print(f'checking space: {current_y} + {box_height} + {para_height} > {max_height}, first_para={first_para}'); if current_y + box_height + para_height > max_height and not first_para:
                            break

                        paragraphs_to_print.pop(0)

                        if first_para:
                            p = tf.paragraphs[0]
                            first_para = False
                        else:
                            p = tf.add_paragraph()

                        if ptype == 'title':
                            run_title = p.add_run()
                            run_title.text      = text
                            run_title.font.bold = True
                            run_title.font.size = Pt(12)
                            p.space_after = Pt(6)
                            if url and url != 'nan':
                                run_title.hyperlink.address    = url
                                run_title.font.color.rgb       = RGBColor(0, 102, 204)
                                run_title.font.underline       = True
                            else:
                                run_title.font.color.rgb = RGBColor(0, 0, 0)
                        else:
                            p.text      = text
                            p.font.size = Pt(11)
                            if is_sub:
                                p.font.bold    = True
                                p.space_before = Pt(12)
                                p.space_after  = Pt(4)
                            elif is_list:
                                p.space_before = Pt(0)
                                p.space_after  = Pt(2)
                            else:
                                p.space_before = Pt(4)
                                p.space_after  = Pt(6)

                        box_height += para_height

                    txBox.height  = Inches(box_height)
                    current_y    += box_height + 0.15

    prs.save(output_path)
    print(f"Generated successfully: {output_path}")


if __name__ == '__main__':
    create_newsletter()
